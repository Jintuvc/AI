from pathlib import Path
import fitz
import pandas as pd
import docx
from pptx import Presentation

class IngestionAgent:
    def parse_document(self, file_path):
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return "\n".join([page.get_text() for page in fitz.open(file_path)])
        elif ext == '.docx':
            return "\n".join([para.text for para in docx.Document(file_path).paragraphs])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == '.csv':
            return pd.read_csv(file_path).to_string()
        elif ext in ['.txt', '.md']:
            return open(file_path, 'r', encoding='utf-8').read()
        else:
            raise ValueError(f"Unsupported format: {ext}")

    def run(self, file_paths):
        docs = [{"filename": Path(p).name, "content": self.parse_document(p)} for p in file_paths]
        return docs
